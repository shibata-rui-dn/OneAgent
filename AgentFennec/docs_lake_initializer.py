from pathlib import Path
import json

# 必要なライブラリのインポート（未インストールの場合はエラーメッセージを表示）
try:
    import openpyxl
except ImportError:
    print("openpyxl がインストールされていません。Excel ファイルの処理は行えません。")

try:
    import docx2txt
except ImportError:
    print("docx2txt がインストールされていません。Word ファイルの処理は行えません。")

try:
    from PyPDF2 import PdfReader
except ImportError:
    print("PyPDF2 がインストールされていません。PDF ファイルの処理は行えません。")

# PowerPoint の処理には python-pptx を使用します（必要な時に関数内でインポート）

def convert_excel_to_text(file_path: Path) -> dict:
    """
    Excel (xlsx, xls, xlsm) ファイルを読み込み、各シートごとのテキストを生成する。
    戻り値は {シート名: テキスト} の辞書形式。
    """
    texts = {}
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        for sheet in wb.worksheets:
            text = f"Sheet: {sheet.title}\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                text += row_text + "\n"
            texts[sheet.title] = text
        return texts
    except Exception as e:
        print(f"Excel ファイル処理エラー ({file_path}): {e}")
        return {}

def convert_word_to_text(file_path: Path) -> list:
    """
    Word (doc, docx, odt) ファイルをテキスト化する。
    ※ docx2txt はページ単位の分割機能を持たないため、フォームフィード文字('\f')で分割を試み、
       分割結果が空の場合は全文を1ページとして返す。
    戻り値は各ページのテキストを要素とするリスト。
    """
    try:
        text = docx2txt.process(str(file_path))
        pages = text.split('\f')
        pages = [p for p in pages if p.strip()]
        if not pages:
            pages = [text]
        return pages
    except Exception as e:
        print(f"Word ファイル処理エラー ({file_path}): {e}")
        return []

def convert_pdf_to_text(file_path: Path) -> list:
    """
    PDF ファイルを読み込み、各ページのテキストを抽出する。
    戻り値は各ページのテキストを要素とするリスト。
    """
    pages_text = []
    try:
        reader = PdfReader(str(file_path))
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                pages_text.append(page_text)
        return pages_text
    except Exception as e:
        print(f"PDF ファイル処理エラー ({file_path}): {e}")
        return []

def convert_ppt_to_text(file_path: Path) -> list:
    """
    PowerPoint (ppt, pptx) ファイルを読み込み、各スライドのテキストを抽出する。
    戻り値は各スライドのテキストを要素とするリスト。
    """
    try:
        from pptx import Presentation
    except ImportError:
        print("python-pptx がインストールされていません。PowerPoint ファイルの処理は行えません。")
        return []
    
    slides_text = []
    try:
        prs = Presentation(str(file_path))
        for i, slide in enumerate(prs.slides, start=1):
            slide_text = f"Slide: {i}\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text += shape.text + "\n"
            slides_text.append(slide_text)
        return slides_text
    except Exception as e:
        print(f"PowerPoint ファイル処理エラー ({file_path}): {e}")
        return []

def convert_file(file_path: Path):
    """
    拡張子に応じてファイルをテキスト化する。
      ・Excel はシートごとに、Word/pdf/PowerPoint はページ/スライドごとに別々の txt ファイル（連番: id.txt）として出力する際の
        戻り値の形式に合わせた変換処理を行う。
    """
    ext = file_path.suffix.lower()
    if ext in [".xlsx", ".xls", ".xlsm"]:
        return convert_excel_to_text(file_path)
    elif ext in [".doc", ".docx", ".odt"]:
        return convert_word_to_text(file_path)
    elif ext == ".pdf":
        return convert_pdf_to_text(file_path)
    elif ext in [".ppt", ".pptx"]:
        return convert_ppt_to_text(file_path)
    else:
        return None

def clear_docs_lake_dir(output_dir: Path):
    """
    docs_lake_dir 内の .txt ファイルと .json ファイルを削除する。
    """
    if not output_dir.exists():
        return
    for file in output_dir.glob("*"):
        if file.suffix in [".txt", ".json"]:
            try:
                file.unlink()
                #print(f"Removed: {file}")
            except Exception as e:
                print(f"削除エラー ({file}): {e}")

def convert_all_files(input_dir: Path, output_dir: Path):
    """
    input_dir 以下を再帰的に探索し、対象ファイル（xlsx, xls, xlsm, doc, docx, odt, pdf, ppt, pptx）をテキスト化する。
    ・Excel はシートごとに、Word/pdf/PowerPoint はページ/スライドごとに別々の txt ファイル（連番: id.txt）として output_dir に保存。
    ・保存時、元ファイルの正規化したパス（resolve()）も mapping に記録し、Excel なら "sheet"、Word/pdf/PowerPoint なら "page" または "slide" のキーも追加する。
    ・さらに、生成された txt ファイルの output_dir 内での相対パスと id を紐づける辞書も生成する。
    戻り値は以下の形式の辞書：
    {
      "id_to_file": { ... },
      "relative_path_to_id": { ... }
    }
    ※進捗状況は progress_callback の代わりに yield で逐次メッセージを返すジェネレーターとして実装しています。
    """
    id_to_file = {}
    relative_path_to_id = {}
    file_id = 1
    abs_output_dir = output_dir.resolve()
    
    # 事前に対象ファイルの数をカウント（出力先ディレクトリは除外）
    valid_extensions = [".xlsx", ".xls", ".xlsm", ".doc", ".docx", ".odt", ".pdf", ".ppt", ".pptx"]
    all_files = [file for file in input_dir.rglob("*")
                 if file.is_file() and file.suffix.lower() in valid_extensions and not (abs_output_dir in file.resolve().parents)]
    total_files = len(all_files)
    processed_files = 0

    for file_path in input_dir.rglob("*"):
        if file_path.is_file():
            try:
                if abs_output_dir in file_path.resolve().parents:
                    continue
            except Exception:
                pass
            if file_path.suffix.lower() in valid_extensions:
                converted = convert_file(file_path)
                if converted:
                    # Excel の場合：シートごとに分割（converted は dict）
                    if file_path.suffix.lower() in [".xlsx", ".xls", ".xlsm"]:
                        for sheet_name, text in converted.items():
                            txt_filename = f"{file_id}.txt"
                            output_file = output_dir / txt_filename
                            try:
                                output_file.write_text(text, encoding="utf-8")
                                #print(f"Converted {file_path} (sheet: {sheet_name}) -> {output_file}")
                                id_to_file[str(file_id)] = {
                                    "original_file": str(file_path.resolve()),
                                    "sheet": sheet_name
                                }
                                # 出力先ディレクトリからの相対パスをキーとして紐づける
                                rel_path = str(output_file.relative_to(output_dir))
                                relative_path_to_id[rel_path] = str(file_id)
                                file_id += 1
                            except Exception as e:
                                print(f"出力エラー ({output_file}): {e}")
                    else:
                        # Word/pdf/PowerPoint の場合：ページまたはスライドごとに分割（converted は list）
                        for page_index, text in enumerate(converted, start=1):
                            txt_filename = f"{file_id}.txt"
                            output_file = output_dir / txt_filename
                            try:
                                output_file.write_text(text, encoding="utf-8")
                                if file_path.suffix.lower() in [".ppt", ".pptx"]:
                                    #print(f"Converted {file_path} (slide: {page_index}) -> {output_file}")
                                    id_to_file[str(file_id)] = {
                                        "original_file": str(file_path.resolve()),
                                        "page": page_index
                                    }
                                else:
                                    print(f"Converted {file_path} (page: {page_index}) -> {output_file}")
                                    id_to_file[str(file_id)] = {
                                        "original_file": str(file_path.resolve()),
                                        "page": page_index
                                    }
                                # 相対パスと id の紐付け
                                rel_path = str(output_file.relative_to(output_dir))
                                relative_path_to_id[rel_path] = str(file_id)
                                file_id += 1
                            except Exception as e:
                                print(f"出力エラー ({output_file}): {e}")
                processed_files += 1
                # 進捗報告（全体の進捗を計算）
                if total_files > 0 and processed_files%5 == 1:
                    percentage = int((processed_files / total_files) * 100)
                    yield "data: " + json.dumps({
                        "status": "progress",
                        "message": "テキスト化実行中",
                        "progress": percentage
                    })+ "\n\n"
    return {
        "id_to_file": id_to_file,
        "relative_path_to_id": relative_path_to_id
    }

def save_mapping(mapping: dict, output_dir: Path):
    """
    mapping 辞書を JSON 形式で output_dir 内の mapping.json に保存する。
    """
    mapping_file = output_dir / "mapping.json"
    try:
        mapping_file.write_text(json.dumps(mapping, ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"Mapping saved: {mapping_file}")
    except Exception as e:
        print(f"Mapping 保存エラー ({mapping_file}): {e}")
